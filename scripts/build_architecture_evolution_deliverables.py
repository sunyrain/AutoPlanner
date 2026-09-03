"""Build the architecture-evolution PDF report and editable PPTX deck."""

from __future__ import annotations

import argparse
import subprocess
from pathlib import Path

import markdown
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOC = ROOT / "docs" / "ARCHITECTURE_EVOLUTION_TIMELINE.md"
HERO = ROOT / "docs" / "assets" / "architecture-evolution" / "architecture-evolution-hero.png"
OUT = ROOT / "docs" / "deliverables"

NAVY = "071526"
NAVY_2 = "0D2136"
PANEL = "102B42"
PANEL_2 = "15374F"
WHITE = "F4F8FB"
MUTED = "9DB1C2"
CYAN = "42D6E8"
TEAL = "2EBFA5"
AMBER = "F5B84B"
CORAL = "F47C70"
BLUE = "66A3FF"
RED = "E35D6A"
GREEN = "45C486"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def build_html() -> Path:
    raw = DOC.read_text(encoding="utf-8")
    body_md = raw.split("\n", 1)[1] if raw.startswith("# ") else raw
    body = markdown.markdown(
        "[TOC]\n\n" + body_md,
        extensions=["tables", "fenced_code", "toc", "sane_lists"],
        extension_configs={"toc": {"title": "目录", "toc_depth": "2-3"}},
    )
    # The HTML lives under docs/deliverables while Markdown asset links are
    # relative to docs/. Use absolute file URIs for reliable headless printing.
    body = body.replace(
        'src="assets/',
        f'src="{(ROOT / "docs" / "assets").as_uri()}/',
    )
    css = f"""
    @page {{ size: A4; margin: 17mm 16mm 17mm 16mm; }}
    @page:first {{ margin: 0; }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; color: #203140; font-family: "Microsoft YaHei", "Noto Sans CJK SC", sans-serif;
            font-size: 10.2pt; line-height: 1.72; background: white; }}
    .cover {{ width: 210mm; height: 297mm; position: relative; overflow: hidden; color: white;
              background: #{NAVY}; page-break-after: always; }}
    .cover img {{ position: absolute; width: 100%; height: 100%; object-fit: cover; opacity: .82; }}
    .cover::after {{ content: ""; position: absolute; inset: 0;
                     background: linear-gradient(90deg, rgba(4,15,29,.96) 0%, rgba(4,15,29,.72) 44%, rgba(4,15,29,.18) 100%); }}
    .cover-copy {{ position: absolute; z-index: 2; left: 20mm; top: 42mm; width: 118mm; }}
    .eyebrow {{ color: #{CYAN}; font-size: 11pt; letter-spacing: 2px; font-weight: 700; }}
    .cover h1 {{ margin: 9mm 0 5mm; font-size: 32pt; line-height: 1.24; letter-spacing: .4px; }}
    .cover .subtitle {{ font-size: 14pt; line-height: 1.7; color: #dce8f0; }}
    .cover .meta {{ position: absolute; z-index: 2; bottom: 22mm; left: 20mm; font-size: 9.5pt; color: #b7cad7; }}
    main {{ max-width: 178mm; margin: 0 auto; }}
    .toc {{ background: #f0f6f8; border-left: 4px solid #{CYAN}; padding: 6mm 8mm; margin: 0 0 9mm; page-break-after: always; }}
    .toc > ul {{ columns: 2; column-gap: 10mm; padding-left: 5mm; }}
    .toc li {{ margin: 1.2mm 0; break-inside: avoid; }}
    .toc a {{ color: #174c67; text-decoration: none; }}
    h1, h2, h3 {{ color: #0a2c43; line-height: 1.34; page-break-after: avoid; }}
    h2 {{ font-size: 18pt; margin: 12mm 0 5mm; padding-bottom: 2mm; border-bottom: 2px solid #d6e4ea; page-break-before: always; }}
    h2:first-of-type {{ page-break-before: auto; }}
    h3 {{ font-size: 13pt; margin: 8mm 0 3mm; color: #14526d; }}
    p {{ margin: 2.5mm 0; orphans: 3; widows: 3; }}
    main p > img {{ display: block; width: 100%; max-height: 112mm; object-fit: contain;
                    margin: 5mm auto 7mm; border: 1px solid #d7e4e9; border-radius: 3mm;
                    page-break-inside: avoid; }}
    ul, ol {{ padding-left: 6mm; }}
    li {{ margin: 1.3mm 0; }}
    blockquote {{ margin: 0 0 8mm; padding: 4mm 6mm; border-left: 4px solid #{AMBER}; background: #fff8e9; color: #455867; }}
    code {{ font-family: Consolas, monospace; font-size: 8.8pt; color: #0b6074; background: #edf5f7; padding: 0 .6mm; }}
    pre {{ background: #f7fbfc; color: #17384d; padding: 5.5mm 6mm; border-radius: 2mm;
           border: 1px solid #c8dde5; border-left: 5px solid #{CYAN};
           font-family: "Microsoft YaHei", Consolas, monospace; font-size: 11pt;
           font-weight: 500; line-height: 1.68; white-space: pre-wrap; page-break-inside: avoid;
           box-shadow: 0 1.5mm 4mm rgba(18, 63, 82, .07); }}
    pre code {{ color: #17384d; background: transparent; padding: 0;
                font-family: "Microsoft YaHei", Consolas, monospace; font-size: inherit; }}
    table {{ width: 100%; border-collapse: collapse; margin: 4mm 0 7mm; font-size: 8.5pt; page-break-inside: avoid; }}
    th {{ background: #{NAVY_2}; color: white; text-align: left; font-weight: 700; padding: 2.4mm; }}
    td {{ border: 1px solid #cfdee5; padding: 2.2mm; vertical-align: top; }}
    tr:nth-child(even) td {{ background: #f4f8fa; }}
    strong {{ color: #103e56; }}
    a {{ color: #087b96; text-decoration: none; }}
    .report-note {{ margin-top: 10mm; padding: 4mm 5mm; background: #edf6f8; border-radius: 2mm; color: #526777; font-size: 8.5pt; }}
    """
    cover = f"""
    <section class="cover">
      <img src="{HERO.as_uri()}" alt="architecture evolution hero">
      <div class="cover-copy">
        <div class="eyebrow">AUTOPLANNER · ARCHITECTURE REVIEW</div>
        <h1>逆合成项目<br>架构演进与理念进步</h1>
        <div class="subtitle">从单步模型与搜索器，到证据优先、最弱环节验收的 Canonical V4<br>设计直觉 · 隐含假设 · 失效信号 · 迁移决策</div>
      </div>
      <div class="meta">深度架构复盘版 · 审阅基线 main@04d9034 · 2026-07-15<br>覆盖 90 个提交 · 7 个阶段 · 7 张结构图 · 代表案例与当前差距</div>
    </section>
    """
    document = f"""<!doctype html><html lang="zh-CN"><head><meta charset="utf-8">
    <title>AutoPlanner 逆合成架构演进时间线</title><style>{css}</style></head>
    <body>{cover}<main>{body}<div class="report-note">生成说明：本报告由仓库 Git 历史、主线文档、架构记录与版本化案例数据交叉整理；未提交工作区修改不计入历史结论。</div></main></body></html>"""
    path = OUT / "autoplanner-architecture-evolution-report.html"
    path.write_text(document, encoding="utf-8")
    return path


def build_pdf(html_path: Path) -> Path:
    chrome = Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe")
    edge = Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe")
    browser = chrome if chrome.exists() else edge
    if not browser.exists():
        raise RuntimeError("Chrome or Edge is required for PDF rendering")
    output = OUT / "AutoPlanner_逆合成架构演进与理念进步_2026-07-15.pdf"
    cmd = [
        str(browser),
        "--headless",
        "--disable-gpu",
        "--no-sandbox",
        "--allow-file-access-from-files",
        "--no-pdf-header-footer",
        f"--print-to-pdf={output}",
        html_path.as_uri(),
    ]
    subprocess.run(
        cmd,
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=120,
    )
    if not output.exists() or output.stat().st_size < 50_000:
        raise RuntimeError("PDF rendering did not produce a valid output")
    return output


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.slide_no = 0

    def add_slide(self, title: str, kicker: str = "ARCHITECTURE EVOLUTION", dark: bool = True):
        slide = self.prs.slides.add_slide(self.blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(NAVY if dark else WHITE)
        self.slide_no += 1
        if kicker:
            self.text(
                slide,
                0.55,
                0.28,
                6.7,
                0.25,
                kicker,
                8.5,
                CYAN if dark else "16748B",
                bold=True,
                spacing=1.4,
            )
        if title:
            self.text(slide, 0.55, 0.58, 12.0, 0.55, title, 24, WHITE if dark else NAVY, bold=True)
        self.line(slide, 0.55, 1.19, 12.75, 1.19, CYAN if dark else "B7D5DE", 1.0)
        return slide

    def text(
        self,
        slide,
        x,
        y,
        w,
        h,
        text,
        size=14,
        color=WHITE,
        bold=False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        font="Microsoft YaHei",
        spacing=1.0,
    ):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return shape

    def rich(self, slide, x, y, w, h, lines, size=12, color=WHITE, bullet=False):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(0.05)
        tf.margin_top = tf.margin_bottom = Pt(0.02)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.space_after = Pt(7)
            p.line_spacing = 1.1
            if bullet:
                p.text = "• " + p.text
            for run in p.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(size)
                run.font.color.rgb = rgb(color)
        return shape

    def rect(self, slide, x, y, w, h, fill=PANEL, line=CYAN, radius=True):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
        return shape

    def line(self, slide, x1, y1, x2, y2, color=CYAN, width=1.5):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = rgb(color)
        line.line.width = Pt(width)
        return line

    def pill(self, slide, x, y, w, text, fill=CYAN, color=NAVY, size=9.5):
        self.rect(slide, x, y, w, 0.34, fill=fill, line=fill)
        self.text(
            slide,
            x,
            y + 0.035,
            w,
            0.22,
            text,
            size,
            color,
            True,
            PP_ALIGN.CENTER,
            MSO_ANCHOR.MIDDLE,
        )

    def footer(self, slide, source="Git history · main@04d9034"):
        self.text(slide, 0.55, 7.16, 10.5, 0.18, source, 7.5, MUTED)
        self.text(
            slide, 12.1, 7.13, 0.65, 0.2, f"{self.slide_no:02d}", 8, CYAN, True, PP_ALIGN.RIGHT
        )

    def card(self, slide, x, y, w, h, title, body, accent=CYAN, metric=None):
        self.rect(slide, x, y, w, h, fill=PANEL, line="244B63")
        self.rect(slide, x, y, 0.06, h, fill=accent, line=accent, radius=False)
        self.text(slide, x + 0.22, y + 0.18, w - 0.4, 0.32, title, 13, WHITE, True)
        if metric:
            self.text(slide, x + 0.22, y + 0.58, w - 0.4, 0.52, metric, 24, accent, True)
            by = y + 1.17
            bh = h - 1.32
        else:
            by = y + 0.62
            bh = h - 0.78
        self.rich(slide, x + 0.22, by, w - 0.42, bh, body, 9.8, MUTED)


def add_cover(d: Deck):
    s = d.prs.slides.add_slide(d.blank)
    d.slide_no += 1
    s.shapes.add_picture(str(HERO), 0, 0, width=d.prs.slide_width, height=d.prs.slide_height)
    overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(7.2), d.prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = rgb("061322")
    overlay.line.fill.background()
    d.text(s, 0.65, 0.62, 5.4, 0.3, "AUTOPLANNER · HISTORY REVIEW", 9, CYAN, True, spacing=1.5)
    d.text(s, 0.65, 1.24, 5.75, 1.55, "逆合成项目\n架构演进与理念进步", 31, WHITE, True)
    d.text(
        s,
        0.68,
        3.12,
        5.55,
        0.9,
        "从单步模型与搜索器，到证据优先、\n最弱环节验收的 Canonical V4",
        15,
        "D6E4ED",
    )
    d.line(s, 0.68, 4.33, 2.35, 4.33, CYAN, 2)
    d.text(s, 0.68, 5.83, 5.3, 0.5, "90 个提交  ·  7 个阶段  ·  5 类代表案例", 11, WHITE, True)
    d.text(s, 0.68, 6.56, 5.4, 0.28, "审阅基线 main@04d9034 · 2026-07-15", 8.5, MUTED)


def build_pptx() -> Path:
    d = Deck()
    add_cover(d)

    s = d.add_slide("结论：进步的核心不是模型更大，而是事实权威更窄")
    d.text(
        s,
        0.62,
        1.48,
        4.3,
        1.25,
        "从“生成一条看似合理的路线”\n到“在硬预算内形成可重放、\n逐边有证据的多样路线组合”",
        20,
        WHITE,
        True,
    )
    d.text(
        s,
        0.62,
        3.05,
        4.1,
        0.75,
        "每次架构升级，都在回答同一个问题：\n谁有资格把 proposal 升级为事实？",
        12,
        MUTED,
    )
    stages = [
        ("模型/搜索器", "局部预测", CORAL),
        ("CascadeBoard", "全局补全", AMBER),
        ("Agentic Blackboard", "协作编排", BLUE),
        ("Evidence-first", "逐边证据", TEAL),
        ("Canonical V4", "单一权威", CYAN),
    ]
    for i, (a, b, c) in enumerate(stages):
        y = 1.47 + i * 0.98
        d.rect(s, 5.25, y, 6.95, 0.72, fill=PANEL_2, line=c)
        d.pill(s, 5.48, y + 0.18, 1.0, f"{i + 1:02d}", c, NAVY)
        d.text(s, 6.68, y + 0.14, 2.35, 0.23, a, 12, WHITE, True)
        d.text(s, 9.26, y + 0.16, 2.55, 0.22, b, 10.5, c, True, PP_ALIGN.RIGHT)
    d.footer(s)

    s = d.add_slide("七个阶段：从模型流水线到陌生目标证据闭环")
    stage_data = [
        ("I", "04-25", "模型 + MCTS", CORAL),
        ("II", "04-26—05-05", "CascadeBoard", AMBER),
        ("III", "05-17—06-08", "Codex-entry", BLUE),
        ("IV", "06-09—06-24", "Agentic Blackboard", "7E8CFF"),
        ("V", "07-10—07-12", "Evidence-first", TEAL),
        ("VI", "07-13", "Canonical V4", CYAN),
        ("VII", "07-14—07-15", "Blind + target-only", GREEN),
    ]
    for i, (n, date, name, c) in enumerate(stage_data):
        row = 0 if i < 4 else 1
        col = i if i < 4 else i - 4
        x = 0.62 + col * 3.12 + (0 if row == 0 else 1.55)
        y = 1.55 + row * 2.42
        d.card(s, x, y, 2.78, 1.82, f"阶段 {n}", [date, name], c)
    d.text(
        s,
        0.72,
        6.45,
        11.8,
        0.35,
        "主线变化：局部能力不断增多，但写入化学事实的路径持续收口。",
        13,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    d.footer(s, "Git log: 2026-04-25—2026-07-15")

    s = d.add_slide("阶段 I—II：从单点指标到路线级全局一致性")
    d.card(
        s,
        0.62,
        1.48,
        3.62,
        4.95,
        "模型流水线 · 2026-04-25",
        [
            "data / expand / conditions / multistep / eval",
            "关注 top-k、GT@5、solve rate、条件 MAE",
            "优点：评估口径与覆盖率审计",
            "局限：条件、库存、来源和拓扑没有统一事实模型",
        ],
        CORAL,
        "49.5%",
    )
    d.text(
        s, 0.92, 3.28, 2.95, 0.42, "酶催化 top-1 仅覆盖 44% 步骤", 9.5, CORAL, True, PP_ALIGN.CENTER
    )
    d.line(s, 4.52, 3.72, 5.38, 3.72, AMBER, 3)
    d.card(
        s,
        5.58,
        1.48,
        3.08,
        4.95,
        "CascadeBoard · 05-04",
        [
            "① Skeleton generation",
            "② Molecular fill",
            "③ Learned route scoring",
            "允许全局补全与回改早期决策",
        ],
        AMBER,
        "3 层",
    )
    d.card(
        s,
        8.92,
        1.48,
        3.78,
        4.95,
        "理念跃迁",
        [
            "单步模型从决策者退为候选生成器 / 感知器",
            "规划对象从“下一步”升级为“整条路线”",
            "但线性 slot chain 仍难表达会聚 AND 语义与证据版本",
        ],
        CYAN,
    )
    d.footer(s, "代表提交 bcdde7e · d7eb58e · f1947b4")

    s = d.add_slide("阶段 III—IV：Codex 获得编排权，确定性验证保留裁决权")
    flow = [
        ("Target", "确定性预检"),
        ("Codex", "选择工作流"),
        ("Typed actions", "安全 / 预算 / 绑定"),
        ("Local tools", "ChemEnzy / 文献 / 视觉"),
        ("Parent proof", "final verdict"),
    ]
    for i, (a, b) in enumerate(flow):
        x = 0.62 + i * 2.52
        d.rect(s, x, 1.58, 2.12, 1.05, fill=PANEL, line=CYAN if i in (1, 4) else "31566C")
        d.text(s, x + 0.12, 1.77, 1.88, 0.26, a, 12, WHITE, True, PP_ALIGN.CENTER)
        d.text(s, x + 0.12, 2.13, 1.88, 0.21, b, 8.6, MUTED, False, PP_ALIGN.CENTER)
        if i < 4:
            d.line(s, x + 2.12, 2.1, x + 2.48, 2.1, CYAN, 2)
    d.card(
        s,
        0.62,
        3.13,
        3.78,
        2.75,
        "Bufotalin 负例",
        [
            "ChemEnzy 可输出许多 raw routes",
            "审计仍可能是 fake_closed_rejected",
            "局部文献锚点不能替代 parent proof",
        ],
        CORAL,
    )
    d.card(
        s,
        4.72,
        3.13,
        3.78,
        2.75,
        "Blackboard 的贡献",
        ["多轮开放研究", "typed summary 与 artifact refs", "预算、失败原因、工具调用可审计"],
        BLUE,
    )
    d.card(
        s,
        8.82,
        3.13,
        3.88,
        2.75,
        "Blackboard 的结构缺陷",
        [
            "协作、扩展、证据和路线投影混存",
            "同一边可能在多个队列状态不一致",
            "适合记录协作，不适合承载 AND/OR 真相",
        ],
        AMBER,
    )
    d.footer(s, "代表提交 15ed2ca · f7d99f8 · 1f3a94c · c1c60a0")

    s = d.add_slide("阶段 V：Evidence-first——广泛生成，狭窄授权")
    planes = [
        ("PROPOSAL", "提出断键与替代", CORAL),
        ("EVIDENCE", "绑定文档与页面", AMBER),
        ("PROOF", "结构 / 映射 / 库存重放", TEAL),
        ("PUBLICATION", "冻结一致 revision", CYAN),
    ]
    for i, (a, b, c) in enumerate(planes):
        x = 0.62 + i * 3.08
        d.card(s, x, 1.48, 2.78, 1.55, a, [b], c)
    d.text(s, 0.62, 3.42, 5.15, 0.42, "Reaction Hypergraph V2", 20, WHITE, True)
    nodes = [
        (1.05, 4.55, CORAL),
        (2.3, 4.0, AMBER),
        (2.3, 5.1, BLUE),
        (3.7, 4.55, TEAL),
        (5.02, 4.55, CYAN),
    ]
    for x, y, c in nodes:
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.42), Inches(0.42))
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(c)
        sh.line.color.rgb = rgb(WHITE)
    for a, b in [
        ((1.47, 4.76), (2.3, 4.21)),
        ((1.47, 4.76), (2.3, 5.31)),
        ((2.72, 4.21), (3.7, 4.76)),
        ((2.72, 5.31), (3.7, 4.76)),
        ((4.12, 4.76), (5.02, 4.76)),
    ]:
        d.line(s, *a, *b, "6D8EA3", 1.5)
    d.text(
        s,
        0.75,
        5.85,
        4.9,
        0.48,
        "完整前体集合 → 一个产物\n保留会聚反应的 AND 语义",
        11,
        MUTED,
        True,
        PP_ALIGN.CENTER,
    )
    d.card(
        s,
        6.12,
        3.4,
        6.58,
        2.82,
        "逐边可信度规则",
        [
            "多个 Codex child = 一个 codex_model 相关组",
            "DOI 字符串或模型转述 ≠ exact evidence",
            "路线强度 = 最弱反应边 + 未闭合叶",
            "Verifier outcome 覆盖 producer 自报状态",
        ],
        CYAN,
    )
    d.footer(s, "代表提交 38e8569 · 4309788 · c613a49")

    s = d.add_slide("阶段 VI：Canonical V4——四个单一权威")
    blocks = [
        ("RunKernel", "事件 · 恢复 · 任务 · 预算", BLUE),
        ("Canonical Hypergraph", "分子 · 超边 · 来源 · 库存", TEAL),
        ("DeficitFrontier", "唯一待办与调度", AMBER),
        ("Proof Portfolio", "最弱证明 · 多样性 · 验收", CYAN),
    ]
    for i, (a, b, c) in enumerate(blocks):
        x = 0.62 + i * 3.08
        d.card(s, x, 1.52, 2.78, 1.62, a, [b], c)
    d.rect(s, 1.34, 3.75, 10.65, 1.62, fill=PANEL_2, line=CYAN)
    d.text(s, 1.68, 4.03, 3.15, 0.34, "GlobalCampaignDirector", 17, WHITE, True)
    d.text(s, 1.68, 4.53, 3.3, 0.38, "路线族 · 共享中间体 · 全局转向", 10, MUTED)
    d.text(s, 5.15, 4.06, 1.65, 0.28, "PROPOSAL ONLY", 9, CORAL, True, PP_ALIGN.CENTER)
    d.line(s, 6.05, 4.5, 7.18, 4.5, CORAL, 2.5)
    d.text(
        s,
        7.38,
        3.98,
        3.9,
        0.62,
        "Host: admission → materialize\n→ validate → evidence → stock",
        12,
        WHITE,
        True,
    )
    d.text(
        s,
        1.0,
        6.15,
        11.3,
        0.36,
        "CLI / API / Web / 恢复 / 导出都是适配器，不再持有第二份化学状态。",
        13,
        CYAN,
        True,
        PP_ALIGN.CENTER,
    )
    d.footer(s, "2026-07-13：13 个提交完成主线权威收口")

    s = d.add_slide("阶段 VII：陌生目标与来源闭环——HTML-first，有界降级")
    chain = [
        ("官方 HTML/XML", GREEN),
        ("PDF 原生文本", TEAL),
        ("本地 OCR", AMBER),
        ("视觉 L0 候选", CORAL),
    ]
    for i, (a, c) in enumerate(chain):
        x = 0.68 + i * 3.12
        d.rect(s, x, 1.58, 2.7, 0.84, fill=PANEL, line=c)
        d.text(s, x + 0.12, 1.83, 2.46, 0.25, a, 12, WHITE, True, PP_ALIGN.CENTER)
        if i < 3:
            d.line(s, x + 2.7, 2.0, x + 3.06, 2.0, c, 2.2)
    d.text(
        s,
        0.72,
        2.68,
        12,
        0.34,
        "上一级已闭合，不进入更昂贵、更不确定的下一级",
        12,
        CYAN,
        True,
        PP_ALIGN.CENTER,
    )
    d.card(
        s,
        0.68,
        3.38,
        3.72,
        2.48,
        "Target-only",
        ["任意陌生 SMILES 冷启动", "全局 campaign + host validation", "失败保留为具名 deficit"],
        BLUE,
    )
    d.card(
        s,
        4.78,
        3.38,
        3.72,
        2.48,
        "Self-evolution",
        ["exact patent row + accepted proof", "原例 replay 后才可入库", "复用边仍从 L0 重新验证"],
        TEAL,
    )
    d.card(
        s,
        8.88,
        3.38,
        3.72,
        2.48,
        "展示边界",
        ["部分条件明确显示为部分", "颜色只投影 proof，不创造 proof", "视觉结果不能直升 L2/L3"],
        AMBER,
    )
    d.footer(s, "代表提交 3dba045 · 1b4afb2 · cc1a48a · 2be8e6c")

    s = d.add_slide("五类代表案例：成功、受限成功与诚实失败必须同时审阅")
    cases = [
        ("Bufotalin", "复杂目标负例", "大量候选仍未闭合", CORAL),
        ("Nirmatrelvir", "确定性 replay", "2 路线 · 0 模型调用", BLUE),
        ("Artemisinin", "精确来源案卷", "2 路线 · 4 库存叶", TEAL),
        ("Paclitaxel", "战略多路线族", "3 路线族仍停在 L0", AMBER),
        ("Blind panel", "陌生目标验收", "4 目标 · 0 假闭合", CYAN),
    ]
    for i, (a, b, c, co) in enumerate(cases):
        x = 0.62 + (i % 3) * 4.08
        y = 1.48 + (i // 3) * 2.2
        d.card(s, x, y, 3.72, 1.82, a, [b, c], co)
    d.text(
        s,
        4.77,
        5.95,
        7.7,
        0.44,
        "审阅重点：系统当前证据允许说到哪一步？",
        16,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    d.footer(s)

    s = d.add_slide("Nirmatrelvir 与 Artemisinin：两种可重放闭环")
    d.card(
        s,
        0.62,
        1.5,
        5.78,
        4.94,
        "Nirmatrelvir · 科学 replay",
        [
            "冻结案卷验证运行内核与恢复一致性",
            "2 条完整路线",
            "12 条规范超边",
            "15 条精确来源记录",
            "7 个库存叶",
            "0 次模型 / 视觉调用",
        ],
        BLUE,
        "2 routes",
    )
    d.card(
        s,
        6.72,
        1.5,
        5.98,
        4.94,
        "Artemisinin · 精确来源案卷",
        [
            "采购边界 A：青蒿酸 → DHAA → 青蒿素",
            "采购边界 B：直接采购 DHAA",
            "2 条验证超边",
            "3 条精确来源记录",
            "4 个库存叶",
            "氢气 / 氧气作为真实叶节点审计",
        ],
        TEAL,
        "2 boundaries",
    )
    d.footer(s, "来源：README.md · config/examples/*_v4_*pack.json")

    s = d.add_slide("Bufotalin 与 Paclitaxel：路线图规模不能替代逐边证明")
    d.card(
        s,
        0.62,
        1.48,
        5.78,
        4.95,
        "Bufotalin · hard negative",
        [
            "ChemEnzy raw routes 很多",
            "审计状态：fake_closed_rejected / partial",
            "文献锚点不等于完整 parent route",
            "推动 Codex-entry、failure critic 与 deterministic verdict",
        ],
        CORAL,
        "≠ solved",
    )
    d.card(
        s,
        6.72,
        1.48,
        5.98,
        4.95,
        "Paclitaxel · bounded showcase",
        [
            "Formal C13 ester：model only",
            "Ojima β-lactam：analogy",
            "Biosynthetic tailoring：analogy",
            "共享 baccatin III 核心",
            "hypotheses_are_not_routes = true",
            "unresolved_is_expected = true",
        ],
        AMBER,
        "3 families · L0",
    )
    d.footer(s, "来源：历史 Bufotalin 审计 · paclitaxel_v4_bounded_showcase_plan.json")

    s = d.add_slide("Blind panel：4 个目标、4 次模型调用、0 条假闭合")
    gates = [
        ("B0\nBlind", 4, GREEN),
        ("B1\nMulti-route", 4, TEAL),
        ("B2\nValidated", 2, BLUE),
        ("B3\nExact evidence", 0, CORAL),
        ("B4\nStock boundary", 3, AMBER),
        ("B5\nAcceptance", 2, CYAN),
    ]
    for i, (label, val, c) in enumerate(gates):
        x = 0.72 + i * 1.62
        maxh = 2.7
        h = max(0.06, maxh * val / 4)
        y = 4.72 - h
        d.rect(s, x, y, 1.06, h, fill=c, line=c, radius=False)
        d.text(s, x, 4.84, 1.06, 0.56, label, 8.3, MUTED, True, PP_ALIGN.CENTER)
        d.text(s, x, y - 0.42, 1.06, 0.3, f"{val}/4", 14, c, True, PP_ALIGN.CENTER)
    d.card(
        s,
        10.2,
        1.48,
        2.5,
        4.35,
        "资源与真实性",
        [
            "4 model calls",
            "72,083 input tokens",
            "25,477 output tokens",
            "151 attempts",
            "43 accepted edges",
            "4/4 within budget",
            "0 false closures",
        ],
        CYAN,
    )
    d.text(
        s,
        0.72,
        5.74,
        9.15,
        0.52,
        "B5 是配置策略 acceptance，不自动等于 evidence-grade 或 procurement-ready。",
        11.5,
        WHITE,
        True,
    )
    d.footer(s, "来源：benchmarks/results/blind_benchmark_summary.v1.json")

    s = d.add_slide("理念进步：六条纵向主线")
    themes = [
        ("智能", "单点预测 → campaign 全局决策", BLUE),
        ("权威", "producer 自报 → 最弱环节验收", CYAN),
        ("状态", "多个投影 → 单一事实源", TEAL),
        ("拓扑", "slot chain → AND/OR hypergraph", AMBER),
        ("证据", "有引用 → 逐边来源生命周期", CORAL),
        ("成本", "搜索参数 → 一等公民合同", GREEN),
    ]
    for i, (a, b, c) in enumerate(themes):
        x = 0.62 + (i % 2) * 6.08
        y = 1.47 + (i // 2) * 1.7
        d.rect(s, x, y, 5.72, 1.34, fill=PANEL, line=c)
        d.pill(s, x + 0.22, y + 0.22, 1.1, a, c, NAVY)
        d.text(s, x + 1.58, y + 0.25, 3.82, 0.62, b, 12, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
    d.footer(s)

    s = d.add_slide("当前成熟度：主线已收口，完整产品证明仍有四个最弱轴")
    items = [
        ("真实库存", "需要版本化供应商快照、offer 过期与撤销", CORAL, "MIGRATING"),
        ("工艺条件", "需完整 condition / procedure proof vector", AMBER, "PARTIAL"),
        ("Blind suite", "integrated 不等于复杂分子 acceptance", BLUE, "IN PROGRESS"),
        ("兼容清理", "需用遥测证明 Blackboard / V3 无权威双写", TEAL, "CONTROLLED"),
    ]
    for i, (a, b, c, state) in enumerate(items):
        x = 0.62 + (i % 2) * 6.08
        y = 1.52 + (i // 2) * 2.15
        d.card(s, x, y, 5.72, 1.72, a, [b], c)
        d.pill(s, x + 4.1, y + 0.18, 1.35, state, c, NAVY, 8)
    d.text(
        s,
        0.85,
        6.06,
        11.62,
        0.52,
        "审阅原则：能力已接入 ≠ 复杂目标已通过；失败样例和成本报告必须与成功样例同等保留。",
        12,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    d.footer(s, "来源：BLACKBOARD_CAPABILITY_MIGRATION.md")

    s = d.add_slide("审阅结论与下一步", kicker="REVIEW TAKEAWAYS")
    d.text(
        s,
        0.72,
        1.48,
        5.0,
        0.72,
        "AutoPlanner 已从“模型集合”\n演进为“权威边界清晰的系统”。",
        22,
        WHITE,
        True,
    )
    d.card(
        s,
        0.72,
        2.72,
        5.0,
        2.7,
        "已经成立",
        [
            "Global Director 只做全局 proposal",
            "四个 canonical 状态权威",
            "逐边证据与最弱环节验收",
            "可恢复、可重放、诚实失败",
        ],
        TEAL,
    )
    d.card(
        s,
        6.12,
        1.48,
        6.58,
        3.94,
        "建议下一个审阅周期固定交付",
        [
            "① 变更前后的唯一写入权威",
            "② Schema / 数据迁移与兼容删除条件",
            "③ 真实 end-to-end artifact + 失败样例",
            "④ 成本、恢复与 zero-model replay",
            "⑤ focused tests 与 blind acceptance 分开报告",
            "⑥ 真实库存与条件覆盖率作为主 KPI",
        ],
        CYAN,
    )
    d.text(
        s,
        6.28,
        5.83,
        6.15,
        0.48,
        "把“架构完成”继续推进为可审计的运行事实。",
        15,
        CYAN,
        True,
        PP_ALIGN.CENTER,
    )
    d.footer(s)

    output = OUT / "AutoPlanner_逆合成架构演进与理念进步_2026-07-15.pptx"
    d.prs.save(output)
    return output


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--skip-pdf", action="store_true")
    args = parser.parse_args()
    OUT.mkdir(parents=True, exist_ok=True)
    html_path = build_html()
    pdf_path = None if args.skip_pdf else build_pdf(html_path)
    pptx_path = build_pptx()
    print(f"html={html_path}")
    if pdf_path:
        print(f"pdf={pdf_path}")
    print(f"pptx={pptx_path}")


if __name__ == "__main__":
    main()
